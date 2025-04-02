from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR_INDEX
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import json
import re
from typing import Dict, List, Any, Optional, Tuple
import os

from models import DataType
from core import Formatting
from utils.logger import setup_console_logger, setup_file_logger

class PPTUpdater:
    def __init__(self, template_path: str, data_path: str):
        self.prs = Presentation(template_path)
        self.logger = setup_console_logger('PPTUpdater')
        self.file_logger = setup_file_logger('PPTUpdater')
        self._load_data(data_path)

    def _load_data(self, data_path: str) -> None:
        error_data = []
        try:
            with open(data_path, 'r', encoding='utf-8') as f:
                json_data = json.load(f)
            
            if 'data' not in json_data:
                self.logger.error("The 'data' key is missing in the JSON file.")
                self.data = []
            else:
                self.data = json_data['data']
            
            self.table_data = []
            self.single_value_data = []
            self.list_data = []
            self.series_data = []
            
            self._categorize_data()

            if error_data:
                self.logger.error(f"Skipping {len(error_data)} items due to incorrect structure.")
                with open('errorData.json', 'w', encoding='utf-8') as json_file:
                    json.dump(error_data, json_file, indent=4)

        except json.JSONDecodeError:
            self.logger.error(f"Failed to parse JSON from file: {data_path}")
            self.data = []
            self.table_data = []
        except IOError:
            self.logger.error(f"Failed to read file: {data_path}")
            self.data = []
            self.table_data = []

    def _categorize_data(self) -> None:
        for item in self.data:
            if isinstance(item, dict):
                data_type = item.get('type')
                if data_type == DataType.TABLE.value:
                    self.table_data.append(item)
                elif data_type == DataType.SINGLE_VALUE.value:
                    self.single_value_data.append(item)
                elif data_type == DataType.LIST.value:
                    self.list_data.append(item)
                elif data_type == DataType.SERIES.value:
                    self.series_data.append(item)
            else:
                self.logger.warning(f"Skipping item due to unexpected structure: {item}")

    def update_table(self, table: Any, table_data: List[Dict]) -> None:
        try:
            required_cols = 0
            for value_item in table_data:
                unique_rows = {item['row'] for item in value_item['values']}
                required_cols = len(value_item['headers'])
            required_rows = len(unique_rows) + 1

            rows = len(table.rows)
            if required_rows > rows:
                self.logger.warning("The table does not have enough rows. Adding new rows.")
                for _ in range(required_rows - rows):
                    table.add_row()
            
            for value_item in table_data:
                for col_idx, header in enumerate(value_item['headers']):
                    table.cell(0, col_idx).text = header
                for value in value_item['values']:
                    row_idx = value['row']
                    col_idx = value['col']-1
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(value['value'])

                    if 'formatting' in value:
                        formatting = Formatting.from_dict(value['formatting'])
                        formatting.apply_to_run(cell.text_frame.paragraphs[0].runs[0])
                    cell_width = Inches(len(cell.text) * 0.2)
                    table.columns[col_idx].width = cell_width
                    
        except Exception as e:
            error_msg = f"Error updating table: {str(e)}"
            self.logger.error(error_msg)

    def update_slides(self) -> None:
        for slide_index in range(len(self.prs.slides)):
            slide = self.prs.slides[slide_index]
            for shape in slide.shapes:
                if shape.has_table:
                    if self.table_data:
                        self.update_table(shape.table, self.table_data)
                        break
                if shape.has_chart:
                    if shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
                        self.update_list_chart(shape)
                    elif shape.chart.chart_type == XL_CHART_TYPE.LINE:
                        self.update_series_line_chart(shape)
                    elif shape.chart.chart_type == XL_CHART_TYPE.PIE_EXPLODED:
                        self.update_pie_chart(slide.shapes)
                        break
            
            self.update_single_value_text(slide_index)

    def update_pie_chart(self, shapes) -> None:
        index = 0
        for shape in shapes:
            if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.PIE_EXPLODED:
                index += 1
                chart = shape.chart
                items = self.series_data if index > 1 else self.list_data
                
                for item in items:
                    if item['id'] in ['MONTHLY_SALES', 'SALES_OVER_TIME']:
                        chart_data = CategoryChartData()
                        if item['id'] == 'MONTHLY_SALES':
                            chart_data.categories = [f"Month {i+1}" for i in range(len(item['values']))]
                        else:
                            chart_data.categories = item['labels']
                        chart_data.add_series(item['displayName'], item['values'])
                        chart.replace_data(chart_data)

    def update_series_line_chart(self, shape) -> None:
        if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.LINE:
            chart = shape.chart
            for item in self.series_data:
                if item['id'] == 'SALES_OVER_TIME':
                    chart_data = CategoryChartData()
                    chart_data.categories = item['labels']
                    chart_data.add_series(item['displayName'], item['values'])
                    chart.replace_data(chart_data)
                    
                    formatting = item.get('formatting')
                    if formatting:
                        for i, fmt in enumerate(formatting):
                            if fmt is not None:
                                point = chart.series[0].points[i]
                                if fmt['bold']:
                                    point.format.fill.solid()
                                    point.format.fill.fore_color.rgb = RGBColor.from_string(fmt['fontColor'])
                                point.format.line.width = Pt(fmt['fontSize'])

                    self.logger.info(f"Updated line chart with data for {item['displayName']}")
                    return
            self.logger.warning("Data for 'SALES_OVER_TIME' not found in series_data.")
        else:
            self.logger.warning("No line chart found or incorrect chart type.")

    def update_list_chart(self, shape) -> None:
        if shape.has_chart and shape.chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED:
            chart = shape.chart
            for item in self.list_data:
                if item['id'] == 'MONTHLY_SALES':
                    chart_data = CategoryChartData()
                    chart_data.categories = [f"Month {i+1}" for i in range(len(item['values']))]
                    chart_data.add_series('Sales', item['values'])
                    chart.replace_data(chart_data)

                    if chart.has_title:
                        chart.chart_title.text_frame.text = item['displayName']
                    
                    try:
                        series = chart.plots[0].series[0]
                        last_point_idx = len(item['values']) - 1
                        point = series.points[last_point_idx]

                        line = point.format.line
                        line.width = Pt(2.5)
                        
                        line.color.rgb = RGBColor(255, 0, 0)

                        self.logger.info("Successfully updated chart formatting")
                    except Exception as e:
                        self.logger.error(f"Error updating chart format: {str(e)}")

                    self.logger.info(f"Updated bar chart with {len(item['values'])} data points")
                    return
            self.logger.warning("MONTHLY_SALES data not found in list_data")
        else:
            self.logger.warning("No chart found or incorrect chart type")

    def update_single_value_text(self, slide_index: int) -> None:
        pattern = r'\{\{([^}]+)\}\}'
        slide = self.prs.slides[slide_index]
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                matches = re.findall(pattern, text)
                
                for match in matches:
                    placeholder_text = f"{{{{{match}}}}}"
                    if match:
                        self._update_text_placeholder(shape, match, placeholder_text, text)

    def _update_text_placeholder(self, shape, match: str, placeholder_text: str, original_text: str) -> None:
        for item in self.single_value_data:
            if item['id'] == match and 'value' in item:
                original_formatting = self._store_original_formatting(shape)
                position = shape.text_frame.paragraphs[0].alignment
                
                new_text = original_text.replace(placeholder_text, str(item['value']))
                shape.text_frame.text = new_text
                
                self._restore_formatting(shape, original_formatting, position)
                
                if 'formatting' in item:
                    formatting = Formatting.from_dict(item['formatting'])
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            formatting.apply_to_run(run)
                return
        
        self.logger.warning(f"No matching data found for placeholder: {placeholder_text}")

    def _store_original_formatting(self, shape) -> List[Dict]:
        original_formatting = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                color = None
                if run.font.color.type == MSO_COLOR_TYPE.RGB:
                    color = run.font.color.rgb
                elif run.font.color.type == MSO_COLOR_TYPE.SCHEME:
                    color = run.font.color.theme_color
                
                original_formatting.append({
                    'font_name': run.font.name,
                    'font_size': run.font.size,
                    'bold': run.font.bold,
                    'italic': run.font.italic,
                    'underline': run.font.underline,
                    'color': color
                })
        return original_formatting

    def _restore_formatting(self, shape, original_formatting: List[Dict], position: Optional[PP_ALIGN] = None) -> None:
        for paragraph in shape.text_frame.paragraphs:
            if position is not None:
                paragraph.alignment = position
            for i, run in enumerate(paragraph.runs):
                if i < len(original_formatting):
                    fmt = original_formatting[i]
                    if fmt['font_name']:
                        run.font.name = fmt['font_name']
                    if fmt['font_size']:
                        run.font.size = fmt['font_size']
                    if fmt['bold'] is not None:
                        run.font.bold = fmt['bold']
                    if fmt['italic'] is not None:
                        run.font.italic = fmt['italic']
                    if fmt['underline'] is not None:
                        run.font.underline = fmt['underline']
                    if fmt['color'] is not None:
                        if isinstance(fmt['color'], RGBColor):
                            run.font.color.rgb = fmt['color']
                        else:
                            run.font.color.theme_color = fmt['color']

    def save(self, output_path: str) -> None:
        try:
            self.prs.save(output_path)
            self.logger.info(f"Presentation saved successfully to {output_path}")
        except Exception as e:
            self.logger.error(f"Failed to save presentation: {str(e)}")
# ppt making